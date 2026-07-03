import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import logging



# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)



def format_worship_together_url(song_name, artist):
    """
    Format a song name and artist into a Worship Together URL format
    Example: 'Goodness of God' by 'Bethel' becomes 'goodness-of-god-bethel'
    If artist is not specified, do not append the artist part.
    """
    def clean_text(text):
        text = text.lower()
        text = re.sub(r'[^a-z0-9\s-]', '', text)
        text = re.sub(r'\s+', '-', text)
        return text.strip('-')
    song_slug = clean_text(song_name)
    artist_slug = clean_text(artist) if artist else ''
    if artist_slug:
        return f"https://www.worshiptogether.com/songs/{song_slug}-{artist_slug}/"
    else:
        return f"https://www.worshiptogether.com/songs/{song_slug}/"


class LyricsFormatter:
    """Class to handle formatting of lyrics and chords for display and PowerPoint generation"""
    @staticmethod
    def parse_lyrics_and_chords(raw_content):
        """
        Parse raw content into structured lyrics and chords
        """
        if not raw_content:
            return {
                'title': 'Unknown Title',
                'sections': []
            }
        # Split content into sections
        sections = re.split(r'\n\s*\n', raw_content)
        parsed_sections = []
        for section in sections:
            if not section.strip():
                continue
            # Check if this is a labeled section
            section_match = re.match(r'^\[?([A-Za-z\s-]+)(\d*)\]?:?\s*(.+)', section, re.IGNORECASE | re.DOTALL)
            if section_match:
                section_type = section_match.group(1).strip()
                section_num = section_match.group(2).strip()
                section_content = section_match.group(3).strip()
                # Parse section content
                lines = section_content.split('\n')
                parsed_lines = []
                for i in range(0, len(lines), 2):
                    if i + 1 < len(lines):
                        chord_line = lines[i].strip()
                        lyric_line = lines[i + 1].strip()
                        if chord_line or lyric_line:
                            parsed_lines.append({
                                'chords': chord_line,
                                'lyrics': lyric_line
                            })
                    else:
                        if lines[i].strip():
                            parsed_lines.append({
                                'chords': '',
                                'lyrics': lines[i].strip()
                            })
                parsed_sections.append({
                    'type': section_type,
                    'number': section_num,
                    'lines': parsed_lines
                })
            else:
                lines = section.strip().split('\n')
                parsed_lines = []
                for i in range(0, len(lines), 2):
                    if i + 1 < len(lines):
                        chord_line = lines[i].strip()
                        lyric_line = lines[i + 1].strip()
                        if chord_line or lyric_line:
                            parsed_lines.append({
                                'chords': chord_line,
                                'lyrics': lyric_line
                            })
                    else:
                        if lines[i].strip():
                            parsed_lines.append({
                                'chords': '',
                                'lyrics': lines[i].strip()
                            })
                if parsed_lines:
                    parsed_sections.append({
                        'type': 'Section',
                        'number': '',
                        'lines': parsed_lines
                    })
        title = 'Unknown Title'
        if sections and sections[0].strip():
            first_section = sections[0].strip().split('\n')
            if len(first_section) >= 1:
                title = first_section[0].strip()
        return {
            'title': title,
            'sections': parsed_sections
        }


def to_snake_case(text):
    text = text.strip().lower()
    text = re.sub(r'[^a-z0-9]+', '_', text)
    text = re.sub(r'_+', '_', text)
    return text.strip('_')


class PowerPointGenerator:
    """Class to generate PowerPoint slides with lyrics and chords"""
    
    def __init__(self):
        self.prs = None
        self.slide = None
        self.title_shape = None
        self.content_shape = None
    
    def create_slide(self, song_data):
        """Create a PowerPoint slide with the song data (adaptive columns and font size, title uses search_name if present, supports transposed keys)"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        slide_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(slide_layout)
        # Add compact title at the very top
        title_left = Inches(0.5)
        title_top = Inches(0.05)
        title_width = Inches(9.0)
        title_height = Inches(0.3)
        title = self.slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )
        title_frame = title.text_frame
        title_frame.clear()
        title_frame.text = (song_data.get('search_name') or song_data.get('title') or 'Unknown Title').strip()
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(16)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 87, 34)
        title_p.alignment = PP_ALIGN.CENTER
        # Content area
        content_top = Inches(0.4)
        content_height = Inches(5.0)
        slide_width = 9.0  # inches
        available_height_pt = 5.0 * 72  # 5.0 inches in points
        min_font_size = 4
        max_font_size = 10
        min_header_size = 6
        max_header_size = 12
        min_columns = 2  # Set minimum columns to 2
        max_columns = 4
        if 'content' in song_data:
            sections = [s for s in song_data['content'].split('\n\n') if s.strip()]
            # Count total lines (including section headers and blank lines)
            total_lines = 0
            section_line_counts = []
            for section in sections:
                lines = [l for l in section.split('\n') if l.strip()]
                section_line_counts.append(len(lines) + 1)  # +1 for blank line after section
                total_lines += len(lines) + 1
            # Try min_columns to max_columns and reduce font size as needed
            found = False
            for num_cols in range(min_columns, max_columns + 1):
                font_size = max_font_size
                header_size = max_header_size
                while font_size >= min_font_size:
                    line_height = font_size + 2  # pt, rough estimate
                    max_lines_per_col = int(available_height_pt / line_height)
                    if total_lines <= num_cols * max_lines_per_col:
                        found = True
                        break
                    font_size -= 1
                    header_size = max(header_size - 1, min_header_size)
                if found:
                    break
            if not found:
                # Use max columns and min font size
                num_cols = max_columns
                font_size = min_font_size
                header_size = min_header_size
            # Distribute sections to columns as contiguous blocks, preserving order
            col_sections = [[] for _ in range(num_cols)]
            # Calculate total lines per section
            total_sections = len(sections)
            sections_per_col = [total_sections // num_cols] * num_cols
            for i in range(total_sections % num_cols):
                sections_per_col[i] += 1
            idx = 0
            for col_idx, count in enumerate(sections_per_col):
                for _ in range(count):
                    col_sections[col_idx].append(sections[idx])
                    idx += 1
            # Improved estimation of actual content height (PowerPoint uses extra spacing)
            header_lines = 0
            content_lines = 0
            spacer_lines = 0
            for col in col_sections:
                first = True
                for section in col:
                    lines = section.split('\n')
                    if not lines or not lines[0].strip():
                        continue
                    header_lines += 1  # Section header
                    content_lines += max(0, len(lines) - 1)  # Content lines
                    if not first:
                        spacer_lines += 1  # Spacer before section
                    first = False
            # Use a 1.25x multiplier for line heights to match PowerPoint's spacing
            header_height = header_lines * (header_size * 1.25)
            content_height_pts = content_lines * (font_size * 1.25)
            spacer_height = spacer_lines * (6 * 1.25)  # Spacer font size is 6pt
            buffer_height = 0.3 * 72  # 0.3 inch buffer
            actual_content_height_pt = header_height + content_height_pts + spacer_height + buffer_height
            # Center content vertically, but ensure it does not overlap the title
            slide_height_pt = 5.625 * 72
            content_top_pt = (slide_height_pt - actual_content_height_pt) / 2
            min_content_margin = 0.1 * 72  # 0.1 inch in points
            min_content_top = (title_top.inches * 72) + (title_height.inches * 72) + min_content_margin
            if content_top_pt < min_content_top:
                content_top_pt = min_content_top
            content_top = Inches(content_top_pt / 72)
            # Layout columns
            col_width = (slide_width - (num_cols - 1) * 0.3) / num_cols
            for i in range(num_cols):
                left = Inches(0.5 + i * (col_width + 0.3))
                col_content = '\n\n'.join(col_sections[i])  # Use double newline to separate sections
                textbox = self.slide.shapes.add_textbox(
                    left, content_top, Inches(col_width), content_height
                )
                self._add_formatted_lyrics(textbox.text_frame, {'content': col_content}, header_size=header_size, content_size=font_size)
        # self._add_disclaimer()  # Removed as per user request
        return self.prs
    
    def _add_formatted_lyrics(self, text_frame, song_data, header_size=12, content_size=10):
        """Add formatted lyrics with chords to the text frame (monospace, bold headers, blank line between sections, adjustable font sizes)"""
        if 'content' in song_data:
            sections = song_data['content'].split('\n\n')
            first_line = True
            for idx, section in enumerate(sections):
                if not section.strip():
                    continue
                lines = section.split('\n')
                if not lines:
                    continue
                # First line is the section header
                section_header = lines[0].strip()
                if first_line:
                    text_frame.text = section_header
                    para = text_frame.paragraphs[0]
                    para.font.bold = True
                    para.font.size = Pt(header_size)
                    para.font.name = 'Consolas'
                    para.font.color.rgb = RGBColor(74, 111, 165)
                    first_line = False
                else:
                    # Add a small space before each new section
                    spacer = text_frame.add_paragraph()
                    spacer.text = ""
                    spacer.font.size = Pt(6)
                    para = text_frame.add_paragraph()
                    para.text = section_header
                    para.font.bold = True
                    para.font.size = Pt(header_size)
                    para.font.name = 'Consolas'
                    para.font.color.rgb = RGBColor(74, 111, 165)
                # The rest are chords/lyrics
                for line in lines[1:]:
                    if not line.strip():
                        continue
                    para = text_frame.add_paragraph()
                    para.text = line
                    para.font.bold = False
                    para.font.size = Pt(content_size)
                    para.font.name = 'Consolas'
                    para.font.color.rgb = RGBColor(51, 51, 51)
    
    def _process_section_content(self, text_frame, content):
        # No longer needed with new formatting logic
        pass
    
    def _add_disclaimer(self):
        """Add disclaimer at the bottom of the slide"""
        disclaimer_left = Inches(0.5)
        disclaimer_top = Inches(5.1)
        disclaimer_width = Inches(9.0)
        disclaimer_height = Inches(0.4)
        
        disclaimer = self.slide.shapes.add_textbox(
            disclaimer_left, disclaimer_top, disclaimer_width, disclaimer_height
        )
        disclaimer_frame = disclaimer.text_frame
        disclaimer_p = disclaimer_frame.add_paragraph()
        disclaimer_p.text = "For personal and educational use only. Please respect copyright laws."
        disclaimer_p.font.size = Pt(8)
        disclaimer_p.font.italic = True
        disclaimer_p.font.color.rgb = RGBColor(150, 150, 150)
        disclaimer_p.alignment = PP_ALIGN.CENTER
    
    def save(self, filename=None, song_data=None):
        """Save the presentation to a file. If filename is not provided, use the song title in snake_case. If filename is not absolute, save in 'saved_slides' directory."""
        if not self.prs:
            raise ValueError("No presentation has been created")
        if not filename:
            # Use song title or search_name in snake_case
            title = None
            if song_data:
                title = song_data.get('search_name') or song_data.get('title')
            if title:
                filename = to_snake_case(title) + '.pptx'
            else:
                import tempfile
                fd, filename = tempfile.mkstemp(suffix='.pptx')
                os.close(fd)
        # Ensure the file is saved in the correct directory if not absolute
        if not os.path.isabs(filename):
            slides_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'saved_slides')
            if not os.path.exists(slides_dir):
                os.makedirs(slides_dir)
            filename = os.path.join(slides_dir, filename)
        self.prs.save(filename)
        return filename


def generate_pptx_from_song_data(song_data, filename=None):
    """
    Generate a PowerPoint file from song data
    Returns the path to the generated file
    """
    try:
        logger.info("\n" + "="*50)
        logger.info("STARTING POWERPOINT GENERATION")
        logger.info("="*50 + "\n")
        
        # Parse the raw content if needed
        if 'content' in song_data and not song_data.get('sections'):
            parsed_data = LyricsFormatter.parse_lyrics_and_chords(song_data['content'])
            # Merge the parsed data with the original song data
            song_data.update(parsed_data)
        
        # Log song data for debugging
        logger.info("\nGenerating PowerPoint for song:")
        logger.info(f"Title: {song_data.get('title', 'Unknown Title')}")
        logger.info(f"Artist: {song_data.get('artist', 'Unknown Artist')}")
        if 'content' in song_data:
            logger.info("\nContent sections:")
            sections = [s for s in song_data['content'].split('\n\n') if s.strip()]
            for i, section in enumerate(sections, 1):
                logger.info(f"\nSection {i}:")
                logger.info(section)
        
        logger.info("\n" + "-"*50)
        logger.info("CREATING POWERPOINT SLIDE")
        logger.info("-"*50 + "\n")
        
        # Generate PowerPoint
        ppt_gen = PowerPointGenerator()
        ppt_gen.create_slide(song_data)
        
        # Save to specified file or to song title in snake_case
        temp_file = ppt_gen.save(filename=filename, song_data=song_data)
        
        logger.info("\n" + "="*50)
        logger.info(f"POWERPOINT GENERATED: {temp_file}")
        logger.info("="*50 + "\n")
        
        return temp_file
    
    except Exception as e:
        logger.error(f"\nERROR GENERATING POWERPOINT: {str(e)}")
        raise
