import os
import json
import uuid
import yaml
from datetime import datetime, timezone
from typing import List, Optional
from src.utils.pptx_generator import generate_pptx_from_song_data
from src.utils.yaml_models import SongYAML
from pptx.util import Inches, Pt
from pptx import Presentation
import logging

SLIDES_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "saved_slides")

if not os.path.exists(SLIDES_DIR):
    os.makedirs(SLIDES_DIR)


def _slide_path(slide_id: str, format: str = "pptx") -> str:
    """Get the path for a slide file in a specific format."""
    extensions = {"pptx": ".pptx", "yaml": ".yaml", "marp": ".marp.md", "html": ".html"}
    ext = extensions.get(format, f".{format}")
    return os.path.join(SLIDES_DIR, f"{slide_id}{ext}")


def _meta_path(slide_id: str) -> str:
    return os.path.join(SLIDES_DIR, f"{slide_id}.json")


def save_slide(song_data: dict, formats: List[str] | None = None) -> dict:
    """
    Save a slide in one or more formats.

    Args:
        song_data: Song data dictionary containing song metadata and content
        formats: List of formats to save ('yaml', 'marp', 'html', 'pptx').
                Defaults to ['pptx'] for backward compatibility.

    Returns:
        Metadata dictionary with structure:
        {
            'id': str,  # UUID
            'title': str,  # Prefers search_name over title field
            'artist': str | None,
            'key': str | None,
            'filenames': dict[str, str],  # Format -> filename mapping
            'created_at': str,  # ISO 8601 UTC timestamp
            'updated_at': str,  # ISO 8601 UTC timestamp
            'bpm': int | None,  # Optional
            'time_signature': str | None,  # Optional
            'ccli_number': str | None  # Optional
        }

    Note:
        - YAML, Marp, and HTML generation currently create placeholder files
        - Full implementation requires Song_Parser, Marp_Generator, and HTML_Renderer
        - Metadata structure complies with Requirement 14.3
    """
    if formats is None:
        formats = ["pptx"]

    slide_id = str(uuid.uuid4())
    meta_path = _meta_path(slide_id)

    # Current timestamp in ISO 8601 UTC format
    timestamp = datetime.now(timezone.utc).isoformat()

    # Build filenames dict for requested formats
    filenames = {}

    # Generate files in requested formats
    if "pptx" in formats:
        pptx_path = _slide_path(slide_id, "pptx")
        generate_pptx_from_song_data(song_data, filename=pptx_path)
        filenames["pptx"] = os.path.basename(pptx_path)

    if "yaml" in formats:
        yaml_path = _slide_path(slide_id, "yaml")
        # TODO: Implement YAML generation when Song_Parser is available
        # For now, create a placeholder file
        with open(yaml_path, "w", encoding="utf-8") as f:
            f.write("# YAML generation not yet implemented\n")
        filenames["yaml"] = os.path.basename(yaml_path)

    if "marp" in formats:
        marp_path = _slide_path(slide_id, "marp")
        # TODO: Implement Marp generation when Marp_Generator is available
        # For now, create a placeholder file
        with open(marp_path, "w", encoding="utf-8") as f:
            f.write("# Marp generation not yet implemented\n")
        filenames["marp"] = os.path.basename(marp_path)

    if "html" in formats:
        html_path = _slide_path(slide_id, "html")
        # TODO: Implement HTML generation when HTML_Renderer is available
        # For now, create a placeholder file
        with open(html_path, "w", encoding="utf-8") as f:
            f.write("<html><body>HTML generation not yet implemented</body></html>\n")
        filenames["html"] = os.path.basename(html_path)

    # Use the user-inputted song name (search_name) as the title if present
    meta = {
        "id": slide_id,
        "title": song_data.get("search_name") or song_data.get("title"),
        "artist": song_data.get("artist"),
        "key": song_data.get("key"),
        "filenames": filenames,
        "created_at": timestamp,
        "updated_at": timestamp,
    }

    # Add optional fields if present in song_data
    if "bpm" in song_data:
        meta["bpm"] = song_data["bpm"]
    if "time_signature" in song_data:
        meta["time_signature"] = song_data["time_signature"]
    if "ccli_number" in song_data:
        meta["ccli_number"] = song_data["ccli_number"]

    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(meta, f, indent=2)

    return meta


def list_slides() -> List[dict]:
    slides = []
    for fname in os.listdir(SLIDES_DIR):
        if fname.endswith(".json"):
            with open(os.path.join(SLIDES_DIR, fname), "r", encoding="utf-8") as f:
                slides.append(json.load(f))
    return slides


def get_slide(slide_id: str) -> Optional[dict]:
    meta_path = _meta_path(slide_id)
    if not os.path.exists(meta_path):
        return None
    with open(meta_path, "r", encoding="utf-8") as f:
        return json.load(f)


def delete_slide(slide_id: str) -> bool:
    """Delete all files associated with a slide."""
    meta_path = _meta_path(slide_id)
    removed = False

    # Try to delete all possible format files
    for format in ["pptx", "yaml", "marp", "html"]:
        file_path = _slide_path(slide_id, format)
        if os.path.exists(file_path):
            os.remove(file_path)
            removed = True

    # Delete metadata file
    if os.path.exists(meta_path):
        os.remove(meta_path)
        removed = True

    return removed


def update_slide(slide_id: str, request_data: dict) -> dict:
    """
    Update a saved slide with modified YAML song data.
    
    Args:
        slide_id: UUID of the slide to update
        request_data: Request body containing:
            - song: Modified song YAML data (dict)
            - formats: Optional list of formats to regenerate (defaults to existing formats)
    
    Returns:
        Updated metadata dictionary
    
    Raises:
        FileNotFoundError: If the slide doesn't exist
        ValueError: If song data is invalid
    
    Requirements: 14.7
    """
    logger = logging.getLogger(__name__)
    
    # Load existing metadata
    existing_meta = get_slide(slide_id)
    if not existing_meta:
        raise FileNotFoundError(f"Slide {slide_id} not found")
    
    # Extract song data from request
    song_data = request_data.get("song")
    if not song_data:
        raise ValueError("Song data is required")
    
    # Determine which formats to regenerate (default to existing formats)
    formats = request_data.get("formats")
    if formats is None:
        # Use existing formats from metadata
        formats = list(existing_meta.get("filenames", {}).keys())
        if not formats:
            formats = ["pptx"]  # Fallback to PPTX if no formats found
    
    # Validate formats parameter
    valid_formats = {"yaml", "marp", "html", "pdf", "pptx"}
    if not isinstance(formats, list):
        raise ValueError("formats must be a list")
    
    invalid_formats = [f for f in formats if f not in valid_formats]
    if invalid_formats:
        raise ValueError(f"Invalid formats: {invalid_formats}")
    
    # Reconstruct SongYAML from dict for YAML/Marp/HTML generation
    from src.utils.preview import _song_from_payload
    from src.utils.marp_generator import generate_marp, MarpOptions
    from src.utils.html_renderer import render_html
    import yaml
    import tempfile
    from dataclasses import asdict
    
    try:
        song = _song_from_payload({"song": song_data})
    except (ValueError, KeyError, TypeError) as e:
        raise ValueError(f"Invalid song data: {e}")
    
    # Current timestamp in ISO 8601 UTC format
    timestamp = datetime.now(timezone.utc).isoformat()
    
    # Build filenames dict for requested formats
    filenames = {}
    
    # Generate files in requested formats
    if "yaml" in formats:
        yaml_path = _slide_path(slide_id, "yaml")
        with open(yaml_path, "w", encoding="utf-8") as f:
            yaml.safe_dump(asdict(song), f, default_flow_style=False, allow_unicode=True, sort_keys=False)
        filenames["yaml"] = os.path.basename(yaml_path)
        logger.info(f"Regenerated YAML for slide {slide_id}")
    
    if "marp" in formats:
        marp_path = _slide_path(slide_id, "marp")
        options = MarpOptions()
        marp_markdown = generate_marp(song, style="practice", options=options)
        with open(marp_path, "w", encoding="utf-8") as f:
            f.write(marp_markdown)
        filenames["marp"] = os.path.basename(marp_path)
        logger.info(f"Regenerated Marp markdown for slide {slide_id}")
    
    if "html" in formats:
        html_path = _slide_path(slide_id, "html")
        # Generate Marp markdown first
        options = MarpOptions()
        marp_markdown = generate_marp(song, style="practice", options=options)
        
        # Render HTML using Marp CLI
        try:
            render_html(marp_markdown, output_path=html_path, timeout=30)
            filenames["html"] = os.path.basename(html_path)
            logger.info(f"Regenerated HTML for slide {slide_id}")
        except Exception as e:
            logger.warning(f"Marp CLI rendering failed for slide {slide_id}: {e}")
            # Create fallback HTML
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(f"<html><body><h1>Rendering Failed</h1><p>{e}</p></body></html>\n")
            filenames["html"] = os.path.basename(html_path)
    
    if "pptx" in formats:
        pptx_path = _slide_path(slide_id, "pptx")
        # Convert SongYAML back to song_data format for PPTX generator
        legacy_song_data = {
            "title": song.title,
            "search_name": song.title,
            "artist": ", ".join(song.authors) if song.authors else None,
            "key": song.target_key,
            "content": _reconstruct_content_from_song(song),
            "ccli_number": song.ccli_number,
            "bpm": song.bpm,
            "time_signature": song.time_signature,
        }
        generate_pptx_from_song_data(legacy_song_data, filename=pptx_path)
        filenames["pptx"] = os.path.basename(pptx_path)
        logger.info(f"Regenerated PPTX for slide {slide_id}")
    
    # Update metadata
    meta = {
        "id": slide_id,
        "title": song.title,
        "artist": ", ".join(song.authors) if song.authors else None,
        "key": song.target_key,
        "filenames": filenames,
        "created_at": existing_meta.get("created_at", timestamp),  # Preserve original creation time
        "updated_at": timestamp,
    }
    
    # Add optional fields if present
    if song.bpm is not None:
        meta["bpm"] = song.bpm
    if song.time_signature:
        meta["time_signature"] = song.time_signature
    if song.ccli_number:
        meta["ccli_number"] = song.ccli_number
    
    # Save updated metadata
    meta_path = _meta_path(slide_id)
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(meta, f, indent=2)
    
    return meta


def _reconstruct_content_from_song(song: SongYAML) -> str:
    """Reconstruct the legacy content format from SongYAML for PPTX generation."""
    from src.utils.chordpro_parser import reconstruct_brackets
    
    content_parts = []
    for section_name in song.arrangement:
        section = song.sections.get(section_name)
        if not section:
            continue
        
        content_parts.append(section_name)
        for line in section.lines:
            content_parts.append(reconstruct_brackets(line))
        content_parts.append("")  # Blank line between sections
    
    return "\n".join(content_parts)


def get_slide_file(slide_id: str, format: str = "pptx") -> Optional[str]:
    """
    Get the path to a slide file in the specified format.

    Args:
        slide_id: UUID of the slide
        format: File format ('pptx', 'yaml', 'marp', or 'html')

    Returns:
        Path to the file if it exists, None otherwise
    """
    file_path = _slide_path(slide_id, format)
    if os.path.exists(file_path):
        return file_path
    return None


def compile_slides_with_index() -> str:
    """
    Compile all saved slides into a single pptx file with an index slide at the front.
    Returns the path to the compiled pptx file.
    """
    logger = logging.getLogger(__name__)
    compiled_path = os.path.join(SLIDES_DIR, "all_songs.pptx")
    slides = list_slides()
    logger.info(f"Found {len(slides)} slides to compile.")
    if not slides:
        logger.error("No slides to compile.")
        raise Exception("No slides to compile.")
    # Order slides alphabetically by title (case-insensitive)
    slides = sorted(slides, key=lambda x: (x.get("title") or "").lower())
    logger.info(f"Slides after sorting: {[s.get('title') for s in slides]}")
    # Start with a new presentation
    prs = Presentation()
    # Set slide size to match individual song slides (10 x 5.625 inches)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    # Add index slide
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Song Index"
    # Add song titles as bullet points
    tf = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(1.5), width=Inches(8), height=Inches(4)
    ).text_frame
    song_slide_indices = []
    for meta in slides:
        p = tf.add_paragraph()
        p.text = f"{meta.get('title', 'Untitled')} - {meta.get('artist', '')}"
        song_slide_indices.append(
            len(prs.slides)
        )  # The next slide to be added will be this song's first slide
    # Remove any leading empty paragraph from the text frame
    while tf.paragraphs and not tf.paragraphs[0].text.strip():
        p = tf.paragraphs[0]
        tf._element.remove(p._p)

    # Append all slides from each pptx, keeping track of the first slide index for each song
    song_first_slide_indices = []
    for idx, meta in enumerate(slides):
        pptx_path = _slide_path(meta["id"], "pptx")
        logger.info(f"Adding slides from: {pptx_path}")
        src_prs = Presentation(pptx_path)
        first_song_slide_idx = len(prs.slides)
        song_first_slide_indices.append(first_song_slide_idx)
        for src_slide in src_prs.slides:
            blank_slide_layout = prs.slide_layouts[6]
            new_slide = prs.slides.add_slide(blank_slide_layout)
            for shape in src_slide.shapes:
                el = shape.element
                new_slide.shapes._spTree.insert_element_before(el, "p:extLst")
    # Add transparent shapes as hyperlinks over each index entry
    index_slide = prs.slides[0]
    textbox_top = 1.5
    textbox_height = 4
    left = Inches(1)
    width = Inches(8)
    tf_paragraphs = tf.paragraphs
    num_paragraphs = len(tf_paragraphs)
    # Calculate the height of each paragraph box
    para_height = textbox_height / num_paragraphs if num_paragraphs else 0.33
    box_height = para_height * 0.75  # 75% of the calculated height
    vertical_offset = (para_height - box_height) / 2  # Center the box on the text line

    for idx, first_song_slide_idx in enumerate(song_first_slide_indices):
        top = Inches(textbox_top + idx * para_height + vertical_offset)
        height = Inches(box_height)
        shape = index_slide.shapes.add_shape(
            1,  # Rectangle
            left,
            top,
            width,
            height,
        )
        shape.fill.solid()
        shape.fill.transparency = 1.0
        shape.line.fill.solid()
        shape.line.fill.transparency = 1.0
        shape.line.width = Pt(0)
        shape.click_action.target_slide = prs.slides[first_song_slide_idx]
        logger.info(
            f"Added clickable area for '{slides[idx].get('title')}' to slide {first_song_slide_idx + 1}"
        )
    prs.save(compiled_path)
    logger.info(f"Compiled PPTX saved to: {compiled_path}")
    return compiled_path


def clear_temp_files():
    """Delete all files in saved_slides except for .pptx and .json files (keeps all saved slides and compiled PPTX)."""
    for fname in os.listdir(SLIDES_DIR):
        if fname.endswith(".pptx") or fname.endswith(".json"):
            continue
        file_path = os.path.join(SLIDES_DIR, fname)
        try:
            os.remove(file_path)
        except Exception as e:
            print(f"Could not remove {file_path}: {e}")
